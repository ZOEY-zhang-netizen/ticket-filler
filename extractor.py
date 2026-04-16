"""Unified text extraction for DOCX / PPTX / PDF / image uploads."""

from __future__ import annotations

import base64
import os
import re
import zipfile
from pathlib import Path


# OOXML content-type fragments used to detect actual file type from ZIP
_CT_WORD = 'wordprocessingml.document'
_CT_PPTX = 'presentationml.presentation'

# Regex to split a table "备注/details" cell on ；before known sub-section markers
# so that 销售时间/使用时间/购买规则 etc. each start their own line.
_CELL_SPLIT_RE = re.compile(
    r'；(?=(?:\d+[、．.）)]\s*)?(?:销售时间|使用时间|适用人群|使用人群|购买规则|退改规则|退改说明|退款规则|退票规则|销售渠道|抖音平台|抖音直播间|直播间|POI)(?:[：:；]|$))'
)
# Identifies a cell as a "details" cell (starts with 销售时间 or numbered 销售时间)
_IS_DETAIL_CELL_RE = re.compile(r'^(?:\d+[、．.）)]\s*)?销售时间[：:]')


def _normalize_text(text: str) -> str:
    text = text.replace('\r', '\n')
    text = re.sub(r'\u3000', ' ', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r' *\n *', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _sniff_ooxml(path: str) -> str | None:
    """Return '.docx' or '.pptx' by inspecting [Content_Types].xml inside the ZIP.
    Returns None if the file is not a recognisable OOXML package."""
    try:
        with zipfile.ZipFile(path, 'r') as zf:
            ct_xml = zf.read('[Content_Types].xml').decode('utf-8', errors='replace')
    except Exception:
        return None
    if _CT_WORD in ct_xml:
        return '.docx'
    if _CT_PPTX in ct_xml:
        return '.pptx'
    return None


def extract_text_from_docx(path: str) -> str:
    """Read DOCX while preserving paragraph and table row boundaries."""
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError as exc:
        raise ImportError('DOCX 解析需要 python-docx') from exc

    doc = Document(path)
    lines: list[str] = []
    for block in doc.element.body:
        if block.tag == qn('w:p'):
            text = ''.join(t.text or '' for t in block.iter(qn('w:t'))).strip()
            if text:
                lines.append(text)
        elif block.tag == qn('w:tbl'):
            for tr in block.iter(qn('w:tr')):
                price_cells: list[str] = []
                detail_lines: list[str] = []
                for tc in tr.findall(qn('w:tc')):
                    cell_text = ''.join(
                        t.text or '' for t in tc.iter(qn('w:t'))
                    ).strip()
                    if not cell_text:
                        continue
                    if _IS_DETAIL_CELL_RE.match(cell_text):
                        # Details cell (备注): split on ；before section markers
                        # so that 销售时间/使用时间/购买规则 each start their own line
                        parts = _CELL_SPLIT_RE.split(cell_text)
                        detail_lines.extend(p.strip() for p in parts if p.strip())
                    else:
                        price_cells.append(cell_text)
                if price_cells:
                    lines.append('\t'.join(price_cells))
                lines.extend(detail_lines)
    return _normalize_text('\n'.join(lines))


def extract_text_from_pptx(path: str) -> str:
    """Read PPTX: extract text from all shapes on all slides in order."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError('PPTX 解析需要 python-pptx') from exc

    prs = Presentation(path)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)
    return _normalize_text('\n'.join(lines))


def extract_text_from_pdf(path: str) -> str:
    try:
        import pdfplumber
    except ImportError as exc:
        raise ImportError('PDF 解析需要 pdfplumber') from exc
    pages: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ''
            if text.strip():
                pages.append(text)
    return _normalize_text('\n\n'.join(pages))


def extract_text_from_image(path: str) -> str:
    """OCR via Anthropic Vision when ANTHROPIC_API_KEY is configured."""
    api_key = os.environ.get('ANTHROPIC_API_KEY')
    if not api_key:
        raise EnvironmentError('图片 OCR 需要 ANTHROPIC_API_KEY 环境变量')
    try:
        from anthropic import Anthropic
    except ImportError as exc:
        raise ImportError('图片 OCR 需要 anthropic') from exc

    media_type = {
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.png': 'image/png',
    }.get(Path(path).suffix.lower(), 'image/jpeg')

    with open(path, 'rb') as fh:
        data = base64.b64encode(fh.read()).decode('utf-8')

    client = Anthropic(api_key=api_key)
    message = client.messages.create(
        model='claude-3-5-sonnet-latest',
        max_tokens=4096,
        messages=[{
            'role': 'user',
            'content': [
                {
                    'type': 'image',
                    'source': {
                        'type': 'base64',
                        'media_type': media_type,
                        'data': data,
                    },
                },
                {
                    'type': 'text',
                    'text': '请完整提取图片中的所有中文表格和正文内容，保留原有顺序、换行和关键字段名称，不要总结。',
                },
            ],
        }],
    )
    parts = [getattr(item, 'text', '') for item in message.content]
    return _normalize_text('\n'.join([p for p in parts if p]))


def extract_text(path: str) -> str:
    suffix = Path(path).suffix.lower()

    # For Office Open XML formats (.docx / .pptx), sniff the actual content
    # type from [Content_Types].xml so that misnamed files (e.g. a .pptx
    # uploaded with a .docx extension) are routed correctly.
    if suffix in {'.docx', '.pptx'}:
        actual = _sniff_ooxml(path) or suffix
        if actual == '.pptx':
            return extract_text_from_pptx(path)
        return extract_text_from_docx(path)

    if suffix == '.pdf':
        return extract_text_from_pdf(path)
    if suffix in {'.jpg', '.jpeg', '.png'}:
        return extract_text_from_image(path)
    raise ValueError(f'不支持的文件格式: {suffix}')
